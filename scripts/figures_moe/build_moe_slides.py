"""Append MoE-extension slides to a COPY of the deck, matching its style.

Reads pretrain_Jiaqi_0527.pptx, adds 6 new slides (Part V divider, experiment
design, and 4 finding pages), renumbers all footers to "/ NN", and writes
pretrain_Jiaqi_0527_MoE.pptx. Original is left untouched.
"""
import copy
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
DECK = HERE.parent.parent / "docs" / "presentation"
SRC = DECK / "pretrain_Jiaqi_0527.pptx"
DST = DECK / "pretrain_Jiaqi_0527_MoE.pptx"
FIGS = DECK / "figures_moe"

# ---- palette (from the deck) ---------------------------------------------
BG = RGBColor(0xF7, 0xF8, 0xFA)
DIVBG = RGBColor(0x1A, 0x25, 0x35)
INK = RGBColor(0x1F, 0x2D, 0x3D)
SUB = RGBColor(0x4A, 0x5A, 0x6B)
BLUE = RGBColor(0x1A, 0x4F, 0x8B)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
GREY = RGBColor(0x6B, 0x7B, 0x8D)
ORANGE = RGBColor(0xC8, 0x84, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF2, 0xF4, 0xF7)
CARD_BLUE = RGBColor(0xE8, 0xEF, 0xF7)
CARD_GREEN = RGBColor(0xE6, 0xF0, 0xEB)
LINE = RGBColor(0xC8, 0xD0, 0xD8)

YAHEI = "Microsoft YaHei"
CONSOLAS = "Consolas"

prs = Presentation(str(SRC))
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _set_bg(slide, color):
    # set solid background fill on the slide
    cSld = slide._element.cSld
    bg = cSld.makeelement(qn("p:bg"), {})
    bgPr = bg.makeelement(qn("p:bgPr"), {})
    fill = bgPr.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    fill.append(clr)
    bgPr.append(fill)
    bgPr.append(bgPr.makeelement(qn("a:effectLst"), {}))
    bg.append(bgPr)
    cSld.insert(0, bg)


def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, bg)
    return s


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=None,
        line_spc=None, bullets=False):
    """runs: list of (text, size, bold, color, font) OR list of paragraphs
    where each paragraph is such a list. If bullets, each paragraph gets a •."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    # normalize: a single paragraph -> wrap
    if runs and isinstance(runs[0], tuple):
        paras = [runs]
    else:
        paras = runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spc:
            p.line_spacing = Pt(line_spc)
        if bullets:
            _add_bullet(p)
        for (text, size, bold, color, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def _add_bullet(p):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    pPr.set("marL", "285750")
    pPr.set("indent", "-285750")
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
    pPr.append(buFont)
    pPr.append(buChar)


def rrect(slide, x, y, w, h, fill, line=None, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def pagenum(slide, n, total):
    txt(slide, 12.20, 7.00, 1.00, 0.40,
        [(f"{n:02d} / {total}", 11, False, GREY, CONSOLAS)],
        align=PP_ALIGN.RIGHT)


def pic(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def title_block(slide, title, label):
    txt(slide, 0.8, 0.5, 9.2, 0.8, [(title, 32, True, INK, YAHEI)], line_spc=38)
    txt(slide, 8.6, 0.66, 4.2, 0.5, [(label, 16, False, BLUE, YAHEI)],
        align=PP_ALIGN.RIGHT)


# small total placeholder; fixed after we know final count
TOTAL = len(prs.slides) + 7

# =========================================================================
# SLIDE A — Part V divider
# =========================================================================
s = add_slide(DIVBG)
txt(s, 1.5, 2.5, 10.3, 0.4, [("PART V", 14, False, GREY, CONSOLAS)],
    align=PP_ALIGN.CENTER, line_spc=16)
txt(s, 1.5, 3.2, 10.3, 1.0, [("MoE Extension", 44, True, WHITE, YAHEI)],
    align=PP_ALIGN.CENTER, line_spc=52)
txt(s, 1.5, 4.4, 10.3, 0.5,
    [("Do the Spectral Laws Hold for Mixture-of-Experts?", 18, False,
      RGBColor(0xB8, 0xC4, 0xD0), YAHEI)], align=PP_ALIGN.CENTER)
pagenum(s, len(prs.slides), TOTAL)

# =========================================================================
# SLIDE B — Experiment design + progress
# =========================================================================
s = add_slide()
title_block(s, "MoE Extension: Experiment Design", "OLMoE-1B-7B")

# left: why + model card
txt(s, 0.8, 1.45, 5.6, 0.4, [("Motivation", 16, True, BLUE, YAHEI)])
txt(s, 0.8, 1.95, 5.6, 1.5, [
    [("Dense spectral laws (SR/$d$, ", 13, False, SUB, YAHEI),
     ("α", 13, False, SUB, YAHEI),
     (") are validated. ", 13, False, SUB, YAHEI)],
    [("Do they transfer to sparse experts, where each token sees only "
      "top-k of N experts?", 13, False, SUB, YAHEI)],
], bullets=True, line_spc=20)

card = rrect(s, 0.8, 3.4, 5.6, 1.65, CARD_BLUE)
txt(s, 1.0, 3.55, 5.2, 1.4, [
    [("Model: ", 13, True, INK, YAHEI),
     ("allenai/OLMoE-1B-7B-0924", 13, False, SUB, CONSOLAS)],
    [("6.9B total · 1.3B active · 64 experts · top-8", 12.5, False, SUB, YAHEI)],
    [("d=2048 · intermediate=1024 · 16 MoE layers", 12.5, False, SUB, YAHEI)],
    [("10 checkpoints, step 5K → 1.22M (≈1.2T tokens)", 12.5, False, SUB, YAHEI)],
], line_spc=19)

txt(s, 0.8, 5.25, 5.6, 0.4, [("What we measure", 16, True, BLUE, YAHEI)])
txt(s, 0.8, 5.72, 5.8, 1.5, [
    [("Per-expert ", 12.5, False, SUB, YAHEI), ("α", 12.5, True, INK, YAHEI),
     (" & SR/$d$ · Router SR/$d$ · EPR (energy", 12.5, False, SUB, YAHEI)],
    [("equipartition) · ", 12.5, False, SUB, YAHEI),
     ("ψ", 12.5, True, INK, YAHEI),
     (" order parameter · cross-expert alignment", 12.5, False, SUB, YAHEI)],
], bullets=True, line_spc=19)

# right: phase roadmap
txt(s, 6.9, 1.45, 5.7, 0.4, [("Phased Roadmap", 16, True, BLUE, YAHEI)])
phases = [
    ("Phase 0  Proof of concept", "OLMoE loads, fused-expert tensors handled", CARD_GREEN, "✓ DONE", GREEN),
    ("Phase 1  Training dynamics", "10 checkpoints measured (57 min, CPU)", CARD_GREEN, "✓ DONE", GREEN),
    ("Phase 2  Cross-model", "Mixtral + Phi-3.5-MoE: α scales with expert width", CARD_GREEN, "✓ DONE", GREEN),
    ("Phase 3  Architecture", "Shared vs routed experts", CARD, "PLANNED", GREY),
]
yy = 1.95
for head, body, fill, badge, badgecol in phases:
    rrect(s, 6.9, yy, 5.7, 1.12, fill)
    txt(s, 7.05, yy + 0.08, 4.0, 0.35, [(head, 13.5, True, INK, YAHEI)])
    txt(s, 7.05, yy + 0.46, 4.3, 0.6, [(body, 11.5, False, SUB, YAHEI)], line_spc=14)
    txt(s, 11.0, yy + 0.10, 1.5, 0.35, [(badge, 11, True, badgecol, CONSOLAS)],
        align=PP_ALIGN.RIGHT)
    yy += 1.24
pagenum(s, len(prs.slides), TOTAL)

# =========================================================================
# Finding pages helper
# =========================================================================
def finding_page(title, label, fig, formula, bullets_list, highlight,
                 hi_fill=CARD_GREEN, hi_col=GREEN):
    s = add_slide()
    title_block(s, title, label)
    pic(s, FIGS / fig, 0.5, 1.5, 6.7, 4.65)
    # right column
    if formula:
        rrect(s, 7.5, 1.65, 5.0, 0.8, CARD)
        txt(s, 7.5, 1.78, 5.0, 0.5, [(formula, 15, True, INK, CONSOLAS)],
            align=PP_ALIGN.CENTER)
        by = 2.75
    else:
        by = 1.7
    txt(s, 7.5, by, 5.1, 2.4, bullets_list, bullets=True, line_spc=21)
    # highlight box bottom-right
    rrect(s, 7.5, 5.15, 5.05, 1.0, hi_fill)
    txt(s, 7.65, 5.30, 4.75, 0.75, highlight, line_spc=18, anchor=MSO_ANCHOR.MIDDLE)
    pagenum(s, len(prs.slides), TOTAL)
    return s


# Finding A — alpha stability
finding_page(
    "Finding A: MoE α Never Reverses",
    "Structural Stability",
    "moe_fig1_alpha_stability.png",
    "",
    [
        [("Expert α stays at ", 13, False, SUB, YAHEI),
         ("1.44–1.46", 13, True, INK, YAHEI),
         (" for the entire run", 13, False, SUB, YAHEI)],
        [("Δα = +0.3% over 1.2M steps — ", 13, False, SUB, YAHEI),
         ("no phase transition", 13, True, GREEN, YAHEI)],
        [("Dense models reverse (OLMo-2-13B: Δα=+2.71);", 13, False, SUB, YAHEI)],
        [("MoE experts are born in the Lévy regime (α<2)", 13, False, SUB, YAHEI)],
    ],
    [[("MoE experts are ", 13, False, INK, YAHEI),
      ("“born solid”", 13, True, GREEN, YAHEI),
      (": structure is fixed at init,", 13, False, INK, YAHEI)],
     [("training only rotates it — α-reversal early warning does not apply.",
       13, False, INK, YAHEI)]],
)

# Finding B — SR/d convergence
finding_page(
    "Finding B: SR/d Obeys the Dense Law",
    "Cross-Architecture Universality",
    "moe_fig2_srd_convergence.png",
    "0.040 + 0.61/√d = 0.0535",
    [
        [("Per-expert SR/$d$ converges to ", 13, False, SUB, YAHEI),
         ("0.052", 13, True, INK, YAHEI)],
        [("within ", 13, False, SUB, YAHEI),
         ("2.3%", 13, True, GREEN, YAHEI),
         (" of the Dense prediction (d=2048)", 13, False, SUB, YAHEI)],
        [("Two phases: compression (→410K) then", 13, False, SUB, YAHEI)],
        [("specialization (slow decline + rising spread)", 13, False, SUB, YAHEI)],
    ],
    [[("The universal compression law is ", 13, False, INK, YAHEI),
      ("architecture-agnostic", 13, True, GREEN, YAHEI),
      (" —", 13, False, INK, YAHEI)],
     [("it holds per-expert, not just for dense matrices.",
       13, False, INK, YAHEI)]],
)

# Finding C — two-phase dynamics
finding_page(
    "Finding C: EPR & ψ Track Specialization",
    "Two-Phase Dynamics",
    "moe_fig3_two_phase.png",
    "",
    [
        [("EPR ", 13, True, INK, YAHEI),
         ("U-curve", 13, True, GREEN, YAHEI),
         (": equilibration → specialization", 13, False, SUB, YAHEI)],
        [("ψ rises +15%: experts sharpen their main", 13, False, SUB, YAHEI)],
        [("singular direction (functional specialization)", 13, False, SUB, YAHEI)],
        [("Router SR/$d$ frozen from step 5K — routing", 13, False, SUB, YAHEI)],
        [("geometry fixed extremely early", 13, False, SUB, YAHEI)],
    ],
    [[("With α flat, ", 13, False, INK, YAHEI),
      ("EPR is the sensitive health signal", 13, True, GREEN, YAHEI),
      ("", 13, False, INK, YAHEI)],
     [("for MoE — monitor it, not α, for expert collapse.",
       13, False, INK, YAHEI)]],
    hi_fill=CARD_BLUE, hi_col=BLUE,
)

# Finding D — MoE vs Dense
finding_page(
    "Finding D: MoE vs Dense Contrast",
    "Summary of Differences",
    "moe_fig4_moe_vs_dense.png",
    "",
    [
        [("Dense: α drops 6.5→3.2 then ", 13, False, SUB, YAHEI),
         ("reverses", 13, True, INK, YAHEI)],
        [("MoE: α flat at 1.46 (below Lévy α=2)", 13, False, SUB, YAHEI)],
        [("Dynamic range: α 0.3% · SR/$d$ 12% · ", 13, False, SUB, YAHEI),
         ("EPR 76%", 13, True, GREEN, YAHEI)],
        [("First per-expert α measurement reported —", 13, False, SUB, YAHEI)],
        [("fills a clear gap in the HT-SR literature", 13, False, SUB, YAHEI)],
    ],
    [[("MoE training is ", 13, False, INK, YAHEI),
      ("a different regime", 13, True, GREEN, YAHEI),
      (": no relaxation,", 13, False, INK, YAHEI)],
     [("ballistic SR/d convergence (β=1.85), info-bottleneck α.",
       13, False, INK, YAHEI)]],
)

# Finding E — Phase 2: expert size sets the alpha regime (3-model staircase)
finding_page(
    "Finding E: Expert Width Sets the α Regime",
    "Phase 2 · Cross-Model",
    "moe_fig5_alpha_vs_expert_size.png",
    "int 1024 → 6400 → 14336  ⇒  α 1.46 → 3.03 → 4.00",
    [
        [("Three models form a clean ", 13, False, SUB, YAHEI),
         ("monotone staircase", 13, True, INK, YAHEI)],
        [("OLMoE int=1024 → ", 13, False, SUB, YAHEI),
         ("α=1.46", 13, True, BLUE, YAHEI),
         ("  (Lévy)", 13, False, SUB, YAHEI)],
        [("Phi-3.5 int=6400 → ", 13, False, SUB, YAHEI),
         ("α=3.03", 13, True, ORANGE, YAHEI),
         ("  (transition)", 13, False, SUB, YAHEI)],
        [("Mixtral int=14336 → ", 13, False, SUB, YAHEI),
         ("α=4.00", 13, True, GREEN, YAHEI),
         ("  (Dense-like)", 13, False, SUB, YAHEI)],
        [("Set by ", 13, False, SUB, YAHEI),
         ("per-expert width", 13, True, INK, YAHEI),
         (", not total model size", 13, False, SUB, YAHEI)],
    ],
    [[("α<2 in MoE is an ", 13, False, INK, YAHEI),
      ("information-bottleneck", 13, True, GREEN, YAHEI),
      (" effect,", 13, False, INK, YAHEI)],
     [("not over-fitting — narrow experts force a heavy tail.",
       13, False, INK, YAHEI)]],
)

# =========================================================================
# Renumber ALL footers to "/ TOTAL"
# =========================================================================
pat = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")
for idx, slide in enumerate(prs.slides, start=1):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        full = sh.text_frame.text.strip()
        if pat.match(full):
            # rewrite first run, keep formatting
            for p in sh.text_frame.paragraphs:
                if p.runs:
                    p.runs[0].text = f"{idx:02d} / {TOTAL}"
                    for extra in p.runs[1:]:
                        extra.text = ""
            break

prs.save(str(DST))
print("Saved", DST, "with", len(prs.slides), "slides (was", len(prs.slides) - 6, ")")
