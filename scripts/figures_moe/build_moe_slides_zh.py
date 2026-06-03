"""Standalone CHINESE version of the 6 MoE-extension slides.

Same layout/colors/figures as build_moe_slides.py, but all text in Chinese.
Builds a fresh 6-slide deck -> pretrain_Jiaqi_0527_MoE_zh.pptx (figures unchanged).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
DECK = HERE.parent.parent / "docs" / "presentation"
DST = DECK / "pretrain_Jiaqi_0527_MoE_zh.pptx"
FIGS = DECK / "figures_moe"

BG = RGBColor(0xF7, 0xF8, 0xFA)
DIVBG = RGBColor(0x1A, 0x25, 0x35)
INK = RGBColor(0x1F, 0x2D, 0x3D)
SUB = RGBColor(0x4A, 0x5A, 0x6B)
BLUE = RGBColor(0x1A, 0x4F, 0x8B)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
GREY = RGBColor(0x6B, 0x7B, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD = RGBColor(0xF2, 0xF4, 0xF7)
CARD_BLUE = RGBColor(0xE8, 0xEF, 0xF7)
CARD_GREEN = RGBColor(0xE6, 0xF0, 0xEB)
ORANGE = RGBColor(0xC8, 0x7A, 0x1A)
SUBWHITE = RGBColor(0xB8, 0xC4, 0xD0)

YAHEI = "Microsoft YaHei"
CONSOLAS = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 10


def _set_bg(slide, color):
    cSld = slide._element.cSld
    bg = cSld.makeelement(qn("p:bg"), {})
    bgPr = bg.makeelement(qn("p:bgPr"), {})
    fill = bgPr.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"),
                           {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    fill.append(clr)
    bgPr.append(fill)
    bgPr.append(bgPr.makeelement(qn("a:effectLst"), {}))
    bg.append(bgPr)
    cSld.insert(0, bg)


def add_slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, bg)
    return s


def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=None,
        line_spc=None, bullets=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    paras = [runs] if (runs and isinstance(runs[0], tuple)) else runs
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spc:
            p.line_spacing = Pt(line_spc)
        if bullets:
            _add_bullet(p)
        for (text, size, bold, color, font) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return tb


def _add_bullet(p):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", "285750")
    pPr.set("indent", "-285750")
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": "•"}))


def rrect(slide, x, y, w, h, fill, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def pagenum(slide, n):
    txt(slide, 12.20, 7.00, 1.00, 0.40,
        [(f"{n:02d} / {TOTAL}", 11, False, GREY, CONSOLAS)], align=PP_ALIGN.RIGHT)


def pic(slide, path, x, y, w, h):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))


def title_block(slide, title, label):
    txt(slide, 0.8, 0.5, 9.6, 0.8, [(title, 30, True, INK, YAHEI)], line_spc=36)
    txt(slide, 8.6, 0.66, 4.2, 0.5, [(label, 16, False, BLUE, YAHEI)],
        align=PP_ALIGN.RIGHT)


# =========================================================================
# 1 — Part divider
# =========================================================================
s = add_slide(DIVBG)
txt(s, 1.5, 2.5, 10.3, 0.4, [("第五部分", 14, False, GREY, YAHEI)],
    align=PP_ALIGN.CENTER, line_spc=16)
txt(s, 1.5, 3.2, 10.3, 1.0, [("MoE 架构扩展", 44, True, WHITE, YAHEI)],
    align=PP_ALIGN.CENTER, line_spc=52)
txt(s, 1.5, 4.4, 10.3, 0.5,
    [("谱学定律在混合专家（MoE）架构上是否依然成立？", 18, False, SUBWHITE, YAHEI)],
    align=PP_ALIGN.CENTER)
pagenum(s, 1)

# =========================================================================
# 2 — Experiment design + progress
# =========================================================================
s = add_slide()
title_block(s, "MoE 扩展：实验设计", "OLMoE-1B-7B")

txt(s, 0.8, 1.45, 5.6, 0.4, [("研究动机", 16, True, BLUE, YAHEI)])
txt(s, 0.8, 1.95, 5.6, 1.5, [
    [("Dense 模型的谱学定律（SR/$d$、", 13, False, SUB, YAHEI),
     ("α", 13, False, SUB, YAHEI),
     ("）已被验证。", 13, False, SUB, YAHEI)],
    [("它们能否迁移到稀疏专家上？——每个 token 仅经过 N 个专家中的 top-k 个",
      13, False, SUB, YAHEI)],
], bullets=True, line_spc=20)

rrect(s, 0.8, 3.4, 5.6, 1.65, CARD_BLUE)
txt(s, 1.0, 3.55, 5.2, 1.4, [
    [("模型：", 13, True, INK, YAHEI),
     ("allenai/OLMoE-1B-7B-0924", 13, False, SUB, CONSOLAS)],
    [("总参数 6.9B · 激活 1.3B · 64 专家 · top-8", 12.5, False, SUB, YAHEI)],
    [("d=2048 · intermediate=1024 · 16 个 MoE 层", 12.5, False, SUB, YAHEI)],
    [("10 个 checkpoint，step 5K → 1.22M（约 1.2T tokens）", 12.5, False, SUB, YAHEI)],
], line_spc=19)

txt(s, 0.8, 5.25, 5.6, 0.4, [("测量指标", 16, True, BLUE, YAHEI)])
txt(s, 0.8, 5.72, 5.8, 1.5, [
    [("逐专家 ", 12.5, False, SUB, YAHEI), ("α", 12.5, True, INK, YAHEI),
     (" 与 SR/$d$ · 路由 SR/$d$ · EPR（能量均分）", 12.5, False, SUB, YAHEI)],
    [("· 序参数 ", 12.5, False, SUB, YAHEI), ("ψ", 12.5, True, INK, YAHEI),
     (" · 跨专家对齐度", 12.5, False, SUB, YAHEI)],
], bullets=True, line_spc=19)

txt(s, 6.9, 1.45, 5.7, 0.4, [("分阶段路线图", 16, True, BLUE, YAHEI)])
phases = [
    ("阶段 0  概念验证", "OLMoE 成功加载，处理 fused 专家张量", CARD_GREEN, "✓ 已完成", GREEN),
    ("阶段 1  训练动力学", "完成 10 个 checkpoint 测量（CPU，57 分钟）", CARD_GREEN, "✓ 已完成", GREEN),
    ("阶段 2  跨模型对比", "Mixtral + Phi-3.5-MoE：α 随专家宽度递增", CARD_GREEN, "✓ 已完成", GREEN),
    ("阶段 3  架构对比", "共享专家 vs 纯路由专家", CARD, "计划中", GREY),
]
yy = 1.95
for head, body, fill, badge, badgecol in phases:
    rrect(s, 6.9, yy, 5.7, 1.12, fill)
    txt(s, 7.05, yy + 0.08, 4.0, 0.35, [(head, 13.5, True, INK, YAHEI)])
    txt(s, 7.05, yy + 0.46, 4.3, 0.6, [(body, 11.5, False, SUB, YAHEI)], line_spc=14)
    txt(s, 11.0, yy + 0.10, 1.5, 0.35, [(badge, 11, True, badgecol, YAHEI)],
        align=PP_ALIGN.RIGHT)
    yy += 1.24
pagenum(s, 2)


# =========================================================================
# Finding helper
# =========================================================================
def finding_page(n, title, label, fig, formula, bullets_list, highlight,
                 hi_fill=CARD_GREEN):
    s = add_slide()
    title_block(s, title, label)
    pic(s, FIGS / fig, 0.5, 1.5, 6.7, 4.65)
    if formula:
        rrect(s, 7.5, 1.65, 5.0, 0.8, CARD)
        txt(s, 7.5, 1.78, 5.0, 0.5, [(formula, 15, True, INK, CONSOLAS)],
            align=PP_ALIGN.CENTER)
        by = 2.75
    else:
        by = 1.7
    txt(s, 7.5, by, 5.1, 2.4, bullets_list, bullets=True, line_spc=21)
    rrect(s, 7.5, 5.15, 5.05, 1.0, hi_fill)
    txt(s, 7.65, 5.30, 4.75, 0.75, highlight, line_spc=18, anchor=MSO_ANCHOR.MIDDLE)
    pagenum(s, n)
    return s


# 发现 A — α 稳定
finding_page(
    3, "发现 A：MoE 的 α 从不反转", "结构稳定性",
    "moe_fig1_alpha_stability.png", "",
    [
        [("专家 α 全程稳定在 ", 13, False, SUB, YAHEI),
         ("1.44–1.46", 13, True, INK, YAHEI)],
        [("120 万步内 Δα = +0.3% —— ", 13, False, SUB, YAHEI),
         ("无相变", 13, True, GREEN, YAHEI)],
        [("Dense 模型会反转（OLMo-2-13B：Δα=+2.71）", 13, False, SUB, YAHEI)],
        [("MoE 专家从初始化起就处于 Lévy 区间（α<2）", 13, False, SUB, YAHEI)],
    ],
    [[("MoE 专家“天生即固态”", 13, True, GREEN, YAHEI),
      ("：结构在初始化时已确定，", 13, False, INK, YAHEI)],
     [("训练只是旋转方向 —— α 反转预警在此不适用。", 13, False, INK, YAHEI)]],
)

# 发现 B — SR/d 收敛
finding_page(
    4, "发现 B：SR/d 遵循 Dense 定律", "跨架构普适性",
    "moe_fig2_srd_convergence.png", "0.040 + 0.61/√d = 0.0535",
    [
        [("逐专家 SR/$d$ 收敛到 ", 13, False, SUB, YAHEI),
         ("0.052", 13, True, INK, YAHEI)],
        [("与 Dense 预测值偏差仅 ", 13, False, SUB, YAHEI),
         ("2.3%", 13, True, GREEN, YAHEI),
         ("（d=2048）", 13, False, SUB, YAHEI)],
        [("两阶段：压缩期（→410K）后进入", 13, False, SUB, YAHEI)],
        [("特化期（缓慢下降 + 专家间分化加大）", 13, False, SUB, YAHEI)],
    ],
    [[("通用压缩定律是 ", 13, False, INK, YAHEI),
      ("跨架构成立的", 13, True, GREEN, YAHEI),
      ("——", 13, False, INK, YAHEI)],
     [("它在逐专家层面成立，而不仅限于稠密矩阵。", 13, False, INK, YAHEI)]],
)

# 发现 C — 两阶段动力学
finding_page(
    5, "发现 C：EPR 与 ψ 刻画特化过程", "两阶段动力学",
    "moe_fig3_two_phase.png", "",
    [
        [("EPR 呈 ", 13, False, SUB, YAHEI),
         ("U 型曲线", 13, True, GREEN, YAHEI),
         ("：均衡化 → 特化", 13, False, SUB, YAHEI)],
        [("ψ 上升 +15%：专家强化各自主奇异方向", 13, False, SUB, YAHEI)],
        [("（即功能特化）", 13, False, SUB, YAHEI)],
        [("路由 SR/$d$ 自 step 5K 起冻结 ——", 13, False, SUB, YAHEI)],
        [("路由几何结构极早固化", 13, False, SUB, YAHEI)],
    ],
    [[("由于 α 不变，", 13, False, INK, YAHEI),
      ("EPR 才是 MoE 的敏感健康信号", 13, True, BLUE, YAHEI)],
     [("—— 监控专家坍缩应看 EPR，而非 α。", 13, False, INK, YAHEI)]],
    hi_fill=CARD_BLUE,
)

# 发现 D — MoE vs Dense
finding_page(
    6, "发现 D：MoE 与 Dense 对比", "差异总览",
    "moe_fig4_moe_vs_dense.png", "",
    [
        [("Dense：α 从 6.5→3.2 后 ", 13, False, SUB, YAHEI),
         ("反转回升", 13, True, INK, YAHEI)],
        [("MoE：α 恒定在 1.46（低于 Lévy α=2）", 13, False, SUB, YAHEI)],
        [("变化幅度：α 0.3% · SR/$d$ 12% · ", 13, False, SUB, YAHEI),
         ("EPR 76%", 13, True, GREEN, YAHEI)],
        [("首次报告逐专家 α 测量 ——", 13, False, SUB, YAHEI)],
        [("填补 HT-SR 文献的明确空白", 13, False, SUB, YAHEI)],
    ],
    [[("MoE 训练是 ", 13, False, INK, YAHEI),
      ("另一种动力学 regime", 13, True, GREEN, YAHEI),
      ("：无弛豫、", 13, False, INK, YAHEI)],
     [("SR/d 弹道式收敛（β=1.85）、α 由信息瓶颈决定。", 13, False, INK, YAHEI)]],
)

# 发现 E — Phase 2：专家宽度决定 α regime（三点阶梯）
finding_page(
    7, "发现 E：专家宽度决定 α regime", "阶段 2 · 跨模型",
    "moe_fig5_alpha_vs_expert_size.png",
    "int 1024 → 6400 → 14336  ⇒  α 1.46 → 3.03 → 4.00",
    [
        [("三个模型构成清晰的 ", 13, False, SUB, YAHEI),
         ("单调阶梯", 13, True, INK, YAHEI)],
        [("OLMoE int=1024 → ", 13, False, SUB, YAHEI),
         ("α=1.46", 13, True, BLUE, YAHEI),
         ("（Lévy）", 13, False, SUB, YAHEI)],
        [("Phi-3.5 int=6400 → ", 13, False, SUB, YAHEI),
         ("α=3.03", 13, True, ORANGE, YAHEI),
         ("（过渡区）", 13, False, SUB, YAHEI)],
        [("Mixtral int=14336 → ", 13, False, SUB, YAHEI),
         ("α=4.00", 13, True, GREEN, YAHEI),
         ("（Dense-like）", 13, False, SUB, YAHEI)],
        [("由 ", 13, False, SUB, YAHEI),
         ("逐专家宽度", 13, True, INK, YAHEI),
         (" 决定，而非模型总规模", 13, False, SUB, YAHEI)],
    ],
    [[("MoE 中的 α<2 是 ", 13, False, INK, YAHEI),
      ("信息瓶颈", 13, True, GREEN, YAHEI),
      (" 效应，", 13, False, INK, YAHEI)],
     [("而非过拟合 —— 窄专家强制形成重尾分布。", 13, False, INK, YAHEI)]],
)

# 发现 F — attn vs FFN：排序一致，但整体压入 Lévy 区
finding_page(
    8, "发现 F：排序一致，整体下移至 Lévy 区", "逐层分解",
    "dense_vs_moe_attn_ffn.png",
    "两者都是 FFN > attn；gap Δα 1.26 → 0.41",
    [
        [("排序保持不变：", 13, False, SUB, YAHEI),
         ("FFN α > attn α", 13, True, INK, YAHEI),
         ("（两者皆然）", 13, False, SUB, YAHEI)],
        [("（Dense：MLP 3.34 > attn 2.08；MoE：1.62 > 1.21）", 13, False, SUB, YAHEI)],
        [("但 MoE 把 ", 13, False, SUB, YAHEI),
         ("每个组件都压到 α<2", 13, True, GREEN, YAHEI)],
        [("且 attn–FFN 的差距缩小（1.26 → 0.41）", 13, False, SUB, YAHEI)],
        [("attn 内部：q/k（α≈1.1）< v/o（α≈1.27）", 13, False, SUB, YAHEI)],
    ],
    [[("MoE 并未反转层次结构 —— 而是把", 13, False, INK, YAHEI),
      ("所有组件", 13, True, GREEN, YAHEI)],
     [("压缩进重尾的 Lévy 区间。", 13, False, INK, YAHEI)]],
    hi_fill=CARD_BLUE,
)

# 发现 G — 专家数量 vs 特化程度
finding_page(
    9, "发现 G：专家越少，特化越强", "逐专家差异",
    "moe_spread_compare.png",
    "OLMoE σ(α)=0.02   vs   Mixtral σ(α)=0.85",
    [
        [("OLMoE（64 专家）：各专家谱性质 ", 13, False, SUB, YAHEI),
         ("近乎相同", 13, True, BLUE, YAHEI)],
        [("—— 细粒度 MoE 趋向同质化", 13, False, SUB, YAHEI)],
        [("Mixtral（8 专家）：α 跨度 2.5–6.5 ——", 13, False, SUB, YAHEI)],
        [("粗粒度 MoE 驱动强烈特化", 13, False, SUB, YAHEI)],
        [("Mixtral FFN α 随深度升高（→6），类似 dense", 13, False, SUB, YAHEI)],
    ],
    [[("专家粒度是一个 ", 13, False, INK, YAHEI),
      ("特化调节旋钮", 13, True, GREEN, YAHEI),
      ("：", 13, False, INK, YAHEI)],
     [("少而宽的专家分化，多而窄的专家趋同。", 13, False, INK, YAHEI)]],
)

# 发现 H — 统一的 α-vs-width（dense + MoE）
finding_page(
    10, "发现 H：矩阵宽度普适地决定 α regime", "Dense + MoE 统一",
    "unified_alpha_vs_width.png",
    "窄矩阵 → Lévy (α<2)   ·   宽矩阵 → α 4–5",
    [
        [("两个族 ", 13, False, SUB, YAHEI),
         ("趋势一致", 13, True, INK, YAHEI),
         ("：越宽 → α 越高", 13, False, SUB, YAHEI)],
        [("Dense Pythia：hidden 512→4096 ⇒ attn α 1.6→5.2", 13, False, SUB, YAHEI)],
        [("MoE：expert int 1024→14336 ⇒ α 1.46→4.00", 13, False, SUB, YAHEI)],
        [("α<2 ", 13, False, SUB, YAHEI),
         ("并非 MoE 独有", 13, True, GREEN, YAHEI),
         ("——小 dense", 13, False, SUB, YAHEI)],
        [("模型的 attention 同样是重尾的", 13, False, SUB, YAHEI)],
    ],
    [[("α regime 由 ", 13, False, INK, YAHEI),
      ("单矩阵宽度", 13, True, GREEN, YAHEI),
      (" 决定，", 13, False, INK, YAHEI)],
     [("而非 MoE/dense 之分 —— 一个统一的调节旋钮。", 13, False, INK, YAHEI)]],
)

prs.save(str(DST))
print("Saved", DST, "with", len(prs.slides), "slides")
