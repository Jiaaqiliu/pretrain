#!/usr/bin/env python3
"""Generate publication-quality PPTX for 'Beyond Loss Curves' presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Constants ─────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors — Academic Journal Palette (Navy + Gold)
C_BG = RGBColor(0xF7, 0xF8, 0xFA)        # 极浅冷灰底
C_DARK = RGBColor(0x1A, 0x25, 0x35)       # 深海军蓝 (section bg)
C_TEXT = RGBColor(0x1F, 0x2D, 0x3D)       # 正文深色
C_SEC = RGBColor(0x4A, 0x5A, 0x6B)        # 二级文字
C_MUTED = RGBColor(0x7C, 0x8A, 0x97)      # 弱化文字
C_ACCENT = RGBColor(0x1A, 0x4F, 0x8B)     # 主强调：学术深蓝
C_TEAL = RGBColor(0x1A, 0x4F, 0x8B)       # 同主蓝 (统一)
C_RED = RGBColor(0x8B, 0x2C, 0x3A)        # 低饱和砖红
C_AMBER = RGBColor(0x96, 0x71, 0x17)      # 暗金/琥珀
C_PURPLE = RGBColor(0x4A, 0x3F, 0x7A)     # 低饱和紫灰
C_EMERALD = RGBColor(0x2D, 0x6A, 0x4F)    # 深松绿
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_BLUE = RGBColor(0xE8, 0xEF, 0xF7) # 极浅蓝灰 (卡片底)
C_LIGHT_GREEN = RGBColor(0xE6, 0xF0, 0xEB) # 极浅绿灰
C_LIGHT_RED = RGBColor(0xF5, 0xEB, 0xED)  # 极浅玫瑰灰
C_LIGHT_AMBER = RGBColor(0xF5, 0xF0, 0xE1) # 极浅暖米色

FONT_MAIN = 'Microsoft YaHei'
FONT_CN = 'Microsoft YaHei'
FONT_MONO = 'Consolas'

FIG_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'paper', 'figures_v2')


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=C_TEXT, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_MAIN, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=C_TEXT, bold_first=False, spacing=1.3, font_name=FONT_MAIN):
    """Add a text box with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.line_spacing = Pt(int(font_size * spacing))
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=15,
                    color=C_SEC, bullet_color=C_ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_MAIN
        p.line_spacing = Pt(int(font_size * 1.6))
        p.level = 0
        # Bullet char
        p.bullet = True
    return txBox


def add_figure(slide, fig_name, left, top, width=None, height=None):
    path = os.path.join(FIG_DIR, fig_name)
    if not os.path.exists(path):
        print(f"  WARNING: {path} not found")
        return None
    if width and height:
        pic = slide.shapes.add_picture(path, left, top, width=width, height=height)
    elif width:
        pic = slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        pic = slide.shapes.add_picture(path, left, top, height=height)
    else:
        pic = slide.shapes.add_picture(path, left, top, height=Inches(4.5))
    return pic


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, text_color=C_TEXT, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(12)
        tf.margin_right = Pt(12)
        tf.margin_top = Pt(8)
        tf.margin_bottom = Pt(8)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = FONT_MAIN
    return shape


def add_page_number(slide, num, total):
    add_text_box(slide, Inches(12.2), Inches(7.0), Inches(1.0), Inches(0.4),
                 f"{num:02d} / {total:02d}", font_size=10, color=C_MUTED,
                 font_name=FONT_MONO, alignment=PP_ALIGN.RIGHT)


# ─── Main ──────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # Blank
    TOTAL = 27

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.2),
                 "Beyond Loss Curves", font_size=52, bold=True, color=C_TEXT,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
                 "Spectral Monitoring for LLM Pretraining", font_size=24,
                 color=C_SEC, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
                 "基于谱结构的大模型预训练监控框架", font_size=18,
                 color=C_MUTED, alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    # Tags
    tags = ["14 Models", "4 Architectures", "~340 Checkpoints", "2 Scales Validated"]
    tag_x_start = Inches(2.8)
    for i, tag in enumerate(tags):
        x = tag_x_start + Inches(i * 2.1)
        color = C_LIGHT_GREEN if i == 3 else C_LIGHT_BLUE
        add_rounded_rect(slide, x, Inches(5.0), Inches(1.9), Inches(0.45),
                        color, tag, font_size=11, text_color=C_ACCENT if i < 3 else C_EMERALD)

    add_text_box(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.4),
                 "组会汇报 · 2026-05-27", font_size=13, color=C_MUTED,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_page_number(slide, 1, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2: The Problem
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '预训练的 "盲飞" 现状', font_size=36, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.5), Inches(0.7), Inches(4.5), Inches(0.5),
                 'The Blind Flight Problem', font_size=18, color=C_ACCENT)

    # Cost table
    table_data = [
        ["Model", "Cost", "GPUs", "Time"],
        ["GPT-4", "~$100M+", "数万 A100", "3-6 月"],
        ["Llama-3-70B", "~$30M", "16K H100", "~3 月"],
        ["DeepSeek-V3", "~$5.5M", "2048 H800", "2 月"],
    ]
    rows, cols = 4, 4
    tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(6.0), Inches(2.5)).table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = table_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            else:
                p.font.color.rgb = C_TEXT

    # Right side - key message
    add_rounded_rect(slide, Inches(7.5), Inches(1.8), Inches(5.0), Inches(2.5),
                    C_DARK, "", border_color=None)
    add_text_box(slide, Inches(8.0), Inches(2.3), Inches(4.0), Inches(0.5),
                 "$100M+ Investment", font_size=22, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(8.0), Inches(3.0), Inches(4.0), Inches(0.8),
                 "Only 1 instrument:\ntraining loss", font_size=15, color=RGBColor(0x8E, 0x9C, 0xAD),
                 alignment=PP_ALIGN.CENTER)

    # Bottom insight
    add_rounded_rect(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.2),
                    C_LIGHT_BLUE, "", border_color=C_ACCENT)
    add_text_box(slide, Inches(1.1), Inches(5.2), Inches(11.0), Inches(0.8),
                 '整个训练过程中，工程师判断"训练是否正常"的核心指标只有一个：training loss',
                 font_size=15, color=C_ACCENT, font_name=FONT_CN)
    add_page_number(slide, 2, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3: Loss Blind Spots
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)

    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 'Loss 的四个根本性盲区', font_size=36, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.5), Inches(0.7), Inches(4.5), Inches(0.5),
                 'Fundamental Blind Spots', font_size=18, color=C_ACCENT)

    blindspots = [
        ("盲区 1", "Loss ↓ ≠ 模型变好", "Memorization 也能降 loss，但泛化结构在退化", C_RED),
        ("盲区 2", "无法跨模型比较", "1B 的 loss=3.0 与 70B 的 loss=3.0 含义完全不同", C_AMBER),
        ("盲区 3", "不知何时降 LR", "Cosine/WSD 都是人为预设，缺乏模型自身信号", C_PURPLE),
        ("盲区 4", "没有早期预警", "结构退化时 loss 仍在下降——直到不可逆", C_TEAL),
    ]
    for i, (label, title, desc, color) in enumerate(blindspots):
        x = Inches(0.8 + i * 3.1)
        add_rounded_rect(slide, x, Inches(1.8), Inches(2.9), Inches(3.2),
                        C_WHITE, "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
        add_text_box(slide, x + Inches(0.2), Inches(2.0), Inches(2.5), Inches(0.4),
                     label, font_size=13, bold=True, color=color, font_name=FONT_CN)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.6),
                     title, font_size=15, bold=True, color=C_TEXT, font_name=FONT_CN)
        add_text_box(slide, x + Inches(0.2), Inches(3.2), Inches(2.5), Inches(1.2),
                     desc, font_size=12, color=C_MUTED, font_name=FONT_CN)

    # Bottom insight
    add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
                    C_LIGHT_AMBER, "")
    add_text_box(slide, Inches(1.1), Inches(5.7), Inches(11.0), Inches(0.6),
                 "核心洞察: 智能不是记住信息，而是发现结构——我们需要直接测量结构",
                 font_size=15, bold=True, color=C_AMBER, font_name=FONT_CN)
    add_page_number(slide, 3, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4: Section - Part I
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_DARK)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.4),
                 "PART I", font_size=14, color=RGBColor(0x6B, 0x7B, 0x8D),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                 "两个核心指标", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.5),
                 "Two Spectral Metrics: SR/d & α", font_size=18,
                 color=RGBColor(0x8E, 0x9C, 0xAD), alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 4, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5: SR/d Definition
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.8),
                 '指标 1: SR/d', font_size=36, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(6.5), Inches(0.7), Inches(5.5), Inches(0.5),
                 'Normalized Stable Rank', font_size=18, color=C_ACCENT)

    # Formula
    add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1.0),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(1.0), Inches(1.75), Inches(5.0), Inches(0.6),
                 "SR(W) = ||W||²F / σ₁²    →    SR/d = SR(W) / d_model",
                 font_size=16, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    # Bullet points
    bullets = [
        "物理含义: 矩阵用了多少百分比的可用维度空间",
        "数学本质: SR = exp(H₂)，即 Rényi-2 谱熵的指数",
        "计算成本: O(mn)，7B 全部层 ~5 min",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(5.5), Inches(2.0),
                   bullets, font_size=14)

    # Insight box
    add_rounded_rect(slide, Inches(0.8), Inches(5.2), Inches(5.5), Inches(1.0),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(1.1), Inches(5.35), Inches(5.0), Inches(0.7),
                 "SR/d ↓ = 谱熵减少 = 信息在压缩\n这不是类比——是数学恒等式",
                 font_size=14, bold=True, color=C_ACCENT, font_name=FONT_CN)

    # Right side - metric boxes
    add_rounded_rect(slide, Inches(7.5), Inches(1.6), Inches(4.5), Inches(1.2),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.5), Inches(1.7), Inches(4.5), Inches(0.7),
                 "~0.40", font_size=36, bold=True, color=C_MUTED,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(7.5), Inches(2.3), Inches(4.5), Inches(0.4),
                 "随机初始化 (Random Init)", font_size=12, color=C_MUTED,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    add_text_box(slide, Inches(7.5), Inches(3.1), Inches(4.5), Inches(0.4),
                 "↓ 训练 ↓", font_size=18, color=C_MUTED,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    add_rounded_rect(slide, Inches(7.5), Inches(3.7), Inches(4.5), Inches(1.2),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.5), Inches(3.8), Inches(4.5), Inches(0.7),
                 "~0.05", font_size=36, bold=True, color=C_TEAL,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(7.5), Inches(4.4), Inches(4.5), Inches(0.4),
                 "训练完成 (Fully Trained)", font_size=12, color=C_MUTED,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    add_text_box(slide, Inches(7.5), Inches(5.2), Inches(4.5), Inches(0.8),
                 "压缩 ~8× → ΔH₂ ≈ -2.0 nats\n(Universal across all models!)",
                 font_size=12, color=C_MUTED, alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 5, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6: Alpha Definition
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.8),
                 '指标 2: α', font_size=36, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(5.5), Inches(0.7), Inches(6.0), Inches(0.5),
                 'Power-Law Exponent', font_size=18, color=C_ACCENT)

    # Formula
    add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.0), Inches(0.9),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(1.0), Inches(1.75), Inches(4.6), Inches(0.5),
                 "P(λ) ~ λ⁻ᵅ    (ESD of WᵀW)",
                 font_size=18, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    bullets = [
        "物理含义: 权重矩阵的内部结构质量",
        "理论基础: Heavy-Tail Self-Regularization (HTSR)",
        "α ↓ = heavy-tail ↑ = 泛化 ↑",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.0), Inches(1.5),
                   bullets, font_size=14)

    # Warning
    add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.0), Inches(1.0),
                    C_LIGHT_RED, "")
    add_text_box(slide, Inches(1.1), Inches(4.95), Inches(4.6), Inches(0.7),
                 "关键发现: α 不是单调下降的——\n在大模型上它会反转 (Reversal)！",
                 font_size=14, bold=True, color=C_RED, font_name=FONT_CN)

    # Right side - alpha ranges table
    ranges_data = [
        ["α 范围", "状态", "含义"],
        ["> 6", "Random", "接近 Marchenko-Pastur"],
        ["4 ~ 6", "Bulk+Spikes", "部分结构化"],
        ["2 ~ 4", "Heavy-Tail", "良好自正则化 ✓"],
        ["< 2", "Over-trained", "Rank collapse 风险"],
    ]
    tbl = slide.shapes.add_table(5, 3, Inches(6.5), Inches(1.6), Inches(5.8), Inches(3.5)).table
    for c in range(3):
        tbl.columns[c].width = Inches(1.9)
    for r in range(5):
        for c in range(3):
            cell = tbl.cell(r, c)
            cell.text = ranges_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            elif r == 1:
                p.font.color.rgb = C_RED
            elif r == 2:
                p.font.color.rgb = C_AMBER
            elif r == 3:
                p.font.color.rgb = C_TEAL
            elif r == 4:
                p.font.color.rgb = C_PURPLE

    add_text_box(slide, Inches(6.5), Inches(5.5), Inches(5.8), Inches(0.8),
                 "Martin & Mahoney (JMLR 2021)\n训练好的网络层遵循 power-law 分布",
                 font_size=12, color=C_MUTED, alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_page_number(slide, 6, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7: Phase Portrait (Hero Figure)
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '两个指标的协同', font_size=36, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.5), Inches(0.7), Inches(4.0), Inches(0.5),
                 'Spectral Phase Portrait', font_size=18, color=C_ACCENT)
    add_figure(slide, 'fig01_phase_portrait.png', Inches(1.5), Inches(1.5), height=Inches(5.5))
    add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.4),
                 "所有模型在 (SR/d, α) 空间中的训练轨迹。小模型收敛到左下角，大模型停滞在高 α 区域",
                 font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_page_number(slide, 7, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8: Section - Part II
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_DARK)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.4),
                 "PART II", font_size=14, color=RGBColor(0x6B, 0x7B, 0x8D),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                 "核心发现", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.5),
                 "Key Findings from 14 Models × 4 Architectures", font_size=18,
                 color=RGBColor(0x8E, 0x9C, 0xAD), alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 8, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9: SR/d Convergence
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8),
                 '发现 1: SR/d 通用收敛', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(9.5), Inches(0.7), Inches(3.5), Inches(0.5),
                 'Universal Convergence', font_size=16, color=C_ACCENT)
    add_figure(slide, 'fig02_srd_convergence.png', Inches(0.5), Inches(1.5), height=Inches(4.8))

    # Right side
    add_rounded_rect(slide, Inches(7.8), Inches(1.8), Inches(4.8), Inches(0.8),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.8), Inches(1.9), Inches(4.8), Inches(0.5),
                 "SR/d(∞) = 0.040 + 0.61 / √d",
                 font_size=16, bold=True, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    bullets = [
        "14 个模型全部收敛到 0.04-0.07",
        "只由层宽度 d 决定，与深度无关",
        "OLMo-2-13B 和 32B (同 d=5120):\n  SR/d 完全相同 (0.043)",
    ]
    add_bullet_list(slide, Inches(7.8), Inches(3.0), Inches(4.8), Inches(2.2), bullets, font_size=13)

    add_rounded_rect(slide, Inches(7.8), Inches(5.5), Inches(4.8), Inches(1.2),
                    C_LIGHT_GREEN, "")
    add_text_box(slide, Inches(8.0), Inches(5.65), Inches(4.4), Inches(0.9),
                 "通用压缩定律: ΔH₂ = -2.04 ± 0.17 nats\n所有 Transformer 精确压缩 ~2 nats 谱熵",
                 font_size=12, bold=True, color=C_EMERALD, font_name=FONT_CN)
    add_page_number(slide, 9, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10: Alpha Reversal
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8),
                 '发现 2: α Reversal', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(9.0), Inches(0.7), Inches(4.0), Inches(0.5),
                 'Structural Degradation Signal', font_size=15, color=C_RED)
    add_figure(slide, 'fig03_alpha_dynamics.png', Inches(0.3), Inches(1.5), width=Inches(7.0))

    # Right text
    add_text_box(slide, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.5),
                 "Loss 下降，但结构在退化！", font_size=18, bold=True, color=C_RED, font_name=FONT_CN)
    bullets = [
        "OLMo-2-13B: α 从 4.25 反转到 6.95 (Δα=+2.71)",
        "5/8 大模型展现 reversal，跨架构一致",
        "MLP 驱动: α_mlp 上升而 α_attn 稳定",
    ]
    add_bullet_list(slide, Inches(7.5), Inches(2.3), Inches(5.0), Inches(2.0), bullets, font_size=13)

    add_rounded_rect(slide, Inches(7.5), Inches(4.5), Inches(5.0), Inches(2.0),
                    RGBColor(0xF2, 0xF4, 0xF7), "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
    add_text_box(slide, Inches(7.7), Inches(4.7), Inches(4.6), Inches(0.3),
                 "直觉 (Intuition)", font_size=11, bold=True, color=C_MUTED, font_name=FONT_MONO)
    add_text_box(slide, Inches(7.7), Inches(5.1), Inches(4.6), Inches(1.2),
                 "嘈杂环境搭积木: 该学的学完后(手没方向), 噪声把精细结构震坏(α↑)。\n降 LR = 减小振动 = 保护结构",
                 font_size=12, color=C_SEC, font_name=FONT_CN)
    add_page_number(slide, 10, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11: Predictive Power
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8),
                 '发现 3: 预测下游性能', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(9.0), Inches(0.7), Inches(4.0), Inches(0.5),
                 'Predictive Power', font_size=16, color=C_ACCENT)
    add_figure(slide, 'fig08_correlation.png', Inches(0.3), Inches(1.4), height=Inches(4.5))

    # Metrics
    metrics = [("−0.90", "Spearman ρ", C_ACCENT), ("0.63", "R² (单变量)", C_TEAL), ("0.84", "R² (+log N)", C_PURPLE)]
    for i, (val, label, color) in enumerate(metrics):
        x = Inches(7.2 + i * 2.0)
        add_text_box(slide, x, Inches(1.6), Inches(1.8), Inches(0.6),
                     val, font_size=28, bold=True, color=color, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.2), Inches(1.8), Inches(0.4),
                     label, font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER, font_name=FONT_CN)

    add_text_box(slide, Inches(7.2), Inches(3.0), Inches(5.5), Inches(0.3),
                 "N=102 (6 models × 17 checkpoints × 5 benchmarks)", font_size=11, color=C_MUTED)

    add_rounded_rect(slide, Inches(7.2), Inches(3.5), Inches(5.5), Inches(0.7),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.2), Inches(3.6), Inches(5.5), Inches(0.4),
                 "Perf ≈ -0.248·log₁₀(SR/d) + 0.062·log₁₀(N) - 0.365",
                 font_size=12, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(7.2), Inches(4.6), Inches(5.5), Inches(0.8),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(7.4), Inches(4.7), Inches(5.1), Inches(0.6),
                 "不需要训练/测试数据或 loss，仅从权重\n就能预测 84% 的下游方差",
                 font_size=13, bold=True, color=C_ACCENT, font_name=FONT_CN)
    add_page_number(slide, 11, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12: Phase Transition
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8),
                 '发现 4: 结构相变', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.5), Inches(0.7), Inches(4.5), Inches(0.5),
                 'Phase Transition at N ≈ 1.7B', font_size=15, color=C_RED)
    add_figure(slide, 'fig07_phase_transition.png', Inches(0.3), Inches(1.5), width=Inches(7.2))

    # Right content
    add_rounded_rect(slide, Inches(7.8), Inches(1.6), Inches(4.8), Inches(1.2),
                    C_WHITE, "", border_color=C_RED)
    add_text_box(slide, Inches(8.0), Inches(1.75), Inches(4.4), Inches(0.9),
                 "< 1.7B: α 能降到 <3 (结构成熟)\n> 1.7B: α stuck >5 (结构不成熟)",
                 font_size=14, color=C_TEXT, font_name=FONT_CN)

    add_rounded_rect(slide, Inches(7.8), Inches(3.1), Inches(4.8), Inches(1.0),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.8), Inches(3.2), Inches(4.8), Inches(0.5),
                 "α(N) ≈ 2.65 + 2.1 × σ((log₁₀N - 9.23) / 0.07)",
                 font_size=11, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.8), Inches(3.7), Inches(4.8), Inches(0.3),
                 "Sigmoid fit R² = 0.97", font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, Inches(7.8), Inches(4.3), Inches(4.8), Inches(1.0),
                    C_LIGHT_AMBER, "")
    add_text_box(slide, Inches(8.0), Inches(4.45), Inches(4.4), Inches(0.7),
                 '"大模型难训" 不只是数据量问题\n——大模型根本性地更难结构化',
                 font_size=13, bold=True, color=C_AMBER, font_name=FONT_CN)
    add_page_number(slide, 12, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13: MLP Bottleneck
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(9.0), Inches(0.8),
                 '发现 5: MLP 是结构瓶颈', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(9.0), Inches(0.7), Inches(4.0), Inches(0.5),
                 'MLP as Structural Bottleneck', font_size=14, color=C_ACCENT)
    add_figure(slide, 'fig06_mlp_attn_gap.png', Inches(0.3), Inches(1.4), height=Inches(4.5))

    # Table
    gap_data = [
        ["Model", "α_attn", "α_mlp", "Gap"],
        ["OLMo-2-1B", "1.35", "3.56", "2.21"],
        ["OLMo-2-32B", "3.44", "7.59", "4.15"],
        ["Mistral-7B", "3.79", "9.22", "5.43"],
    ]
    tbl = slide.shapes.add_table(4, 4, Inches(7.3), Inches(1.8), Inches(5.3), Inches(2.5)).table
    for r in range(4):
        for c in range(4):
            cell = tbl.cell(r, c)
            cell.text = gap_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT_MONO if r > 0 and c > 0 else FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC

    add_rounded_rect(slide, Inches(7.3), Inches(4.8), Inches(5.3), Inches(0.8),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(7.5), Inches(4.9), Inches(5.0), Inches(0.6),
                 "Attention 已结构化 (α<4)，MLP 仍随机 (α>7)",
                 font_size=13, bold=True, color=C_ACCENT, font_name=FONT_CN)
    add_text_box(slide, Inches(7.3), Inches(5.8), Inches(5.3), Inches(0.5),
                 "→ 优先增加 MLP 容量 (更宽 FFN / MoE)",
                 font_size=13, color=C_SEC, font_name=FONT_CN)
    add_page_number(slide, 13, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14: Section - Part III
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_DARK)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.4),
                 "PART III", font_size=14, color=RGBColor(0x6B, 0x7B, 0x8D),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                 "从诊断到行动", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.5),
                 "α-Guided Adaptive LR Schedule", font_size=18,
                 color=RGBColor(0x8E, 0x9C, 0xAD), alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 14, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 15: Alpha-Guided Schedule
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(12.0), Inches(0.8),
                 'α-Guided Schedule', font_size=32, bold=True, color=C_TEXT)
    add_text_box(slide, Inches(8.5), Inches(0.7), Inches(4.5), Inches(0.5),
                 'Data-Driven LR Decay', font_size=16, color=C_ACCENT)

    # Flow steps
    steps = [
        ("每 500 步测量 α (2 层 SVD, ~5s)", C_LIGHT_BLUE, C_ACCENT),
        ("检测: dα/dt > 0 连续 3+ 次?", RGBColor(0xF6,0xF8,0xFB), C_TEXT),
        ("触发 LR Decay (linear → η_min)", C_LIGHT_RED, C_RED),
    ]
    for i, (text, bg, tc) in enumerate(steps):
        y = Inches(1.8 + i * 1.3)
        add_rounded_rect(slide, Inches(1.0), y, Inches(4.5), Inches(0.7), bg, "")
        add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(4.1), Inches(0.5),
                     text, font_size=13, color=tc, font_name=FONT_CN, alignment=PP_ALIGN.CENTER)
        if i < 2:
            add_text_box(slide, Inches(3.0), y + Inches(0.75), Inches(0.5), Inches(0.4),
                         "↓", font_size=16, color=C_MUTED, alignment=PP_ALIGN.CENTER)

    # Right side
    add_text_box(slide, Inches(6.5), Inches(1.6), Inches(5.5), Inches(0.4),
                 "核心优势", font_size=18, bold=True, color=C_TEXT, font_name=FONT_CN)
    bullets = [
        "自动适应不同规模: 410M→83%, 1B→74%",
        "自动适应不同数据: 无需人工调参",
        "信号来自模型自身结构",
    ]
    add_bullet_list(slide, Inches(6.5), Inches(2.1), Inches(5.5), Inches(1.8), bullets, font_size=14)

    add_rounded_rect(slide, Inches(6.5), Inches(4.2), Inches(5.5), Inches(2.2),
                    RGBColor(0xF2, 0xF4, 0xF7), "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
    add_text_box(slide, Inches(6.7), Inches(4.4), Inches(5.0), Inches(0.3),
                 "理论依据 (Theoretical Basis)", font_size=11, bold=True, color=C_MUTED, font_name=FONT_MONO)
    add_text_box(slide, Inches(6.7), Inches(4.8), Inches(5.0), Inches(1.4),
                 "α reversal = SGD noise > spectral gap\n降 η → noise ∝ η² 下降 → 恢复 SNR > 1\n(Wedin's sin-Θ theorem bounds subspace rotation)",
                 font_size=12, color=C_SEC, font_name=FONT_CN)
    add_page_number(slide, 15, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 16: 410M Results
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(12.0), Inches(0.8),
                 '实验验证: 410M 3-Way', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'FineWeb-Edu, 10B tokens, 2 seeds', font_size=14, color=C_MUTED)
    add_figure(slide, 'fig04_3way_comparison.png', Inches(0.3), Inches(1.4), width=Inches(7.0))

    # Results table
    res_data = [
        ["Schedule", "Loss", "α", "Avg ↑"],
        ["Cosine", "2.931", "2.68", "0.459"],
        ["WSD", "2.874", "2.47", "0.467"],
        ["α-Guided", "2.877", "2.44", "0.468"],
    ]
    tbl = slide.shapes.add_table(4, 4, Inches(7.5), Inches(1.8), Inches(5.0), Inches(2.2)).table
    for r in range(4):
        for c in range(4):
            cell = tbl.cell(r, c)
            cell.text = res_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = FONT_MONO if r > 0 and c > 0 else FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            elif r == 3:
                p.font.bold = True

    # Metric boxes
    metric_items = [("+1.95%", "vs Cosine", C_ACCENT), ("+6.3%", "LAMBADA ↑", C_TEAL), ("≈ WSD", "自动匹配", C_AMBER)]
    for i, (val, lab, col) in enumerate(metric_items):
        x = Inches(7.5 + i * 1.7)
        add_rounded_rect(slide, x, Inches(4.5), Inches(1.5), Inches(1.2),
                        RGBColor(0xF2, 0xF4, 0xF7), "")
        add_text_box(slide, x, Inches(4.6), Inches(1.5), Inches(0.5),
                     val, font_size=18, bold=True, color=col, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(5.1), Inches(1.5), Inches(0.4),
                     lab, font_size=10, color=C_MUTED, alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_page_number(slide, 16, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 17: 1B Scale-Up
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(12.0), Inches(0.8),
                 'Scale-Up: 1B 验证', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Advantage Amplifies with Scale', font_size=15, color=C_ACCENT)

    # 1B table
    res1b = [
        ["Schedule", "Loss", "α", "SR/d", "Avg ↑"],
        ["Cosine", "2.739", "4.05", "0.056", "0.486"],
        ["WSD", "2.701", "4.01", "0.054", "0.494"],
        ["α-Guided", "2.709", "3.87", "0.054", "0.498"],
    ]
    tbl = slide.shapes.add_table(4, 5, Inches(0.8), Inches(1.6), Inches(6.0), Inches(2.4)).table
    for r in range(4):
        for c in range(5):
            cell = tbl.cell(r, c)
            cell.text = res1b[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT_MONO if r > 0 and c > 0 else FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            elif r == 3:
                p.font.bold = True

    add_rounded_rect(slide, Inches(0.8), Inches(4.3), Inches(6.0), Inches(0.7),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(1.0), Inches(4.4), Inches(5.6), Inches(0.5),
                 "α-Guided 首次明确超越 WSD (+0.98%)，而非仅持平",
                 font_size=13, bold=True, color=C_ACCENT, font_name=FONT_CN)

    # Cross-scale comparison
    add_text_box(slide, Inches(7.5), Inches(1.5), Inches(5.0), Inches(0.4),
                 "跨规模趋势", font_size=18, bold=True, color=C_TEXT, font_name=FONT_CN)
    comp_data = [
        ["对比", "410M", "1B"],
        ["α-Guided vs Cosine", "+1.95%", "+2.56%"],
        ["α-Guided vs WSD", "+0.11%", "+0.98%"],
        ["Decay point", "83%", "74%"],
        ["LAMBADA ↑", "+6.3%", "+10.2%"],
    ]
    tbl2 = slide.shapes.add_table(5, 3, Inches(7.5), Inches(2.0), Inches(5.0), Inches(3.0)).table
    for r in range(5):
        for c in range(3):
            cell = tbl2.cell(r, c)
            cell.text = comp_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            elif c == 2:
                p.font.bold = True
                p.font.color.rgb = C_ACCENT

    add_rounded_rect(slide, Inches(7.5), Inches(5.5), Inches(5.0), Inches(0.7),
                    C_LIGHT_GREEN, "")
    add_text_box(slide, Inches(7.7), Inches(5.6), Inches(4.6), Inches(0.5),
                 "优势随规模放大: 预计 7B+ 更显著",
                 font_size=13, bold=True, color=C_EMERALD, font_name=FONT_CN)
    add_page_number(slide, 17, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 18: Cross-Architecture
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '跨架构验证', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Cross-Architecture Validation', font_size=16, color=C_ACCENT)
    add_figure(slide, 'fig11_cross_arch.png', Inches(0.3), Inches(1.4), height=Inches(4.2))
    add_figure(slide, 'fig05_srd_vs_d.png', Inches(7.0), Inches(1.4), height=Inches(4.2))

    # Architecture labels
    archs = [("GPT-NeoX (Pythia)", "6 scales, 150 ckpts"),
             ("OLMo-2", "4 scales, 100 ckpts"),
             ("Mistral (GQA)", "v0.1 + v0.3, Δα < 0.02")]
    for i, (name, desc) in enumerate(archs):
        x = Inches(0.8 + i * 4.2)
        add_rounded_rect(slide, x, Inches(6.0), Inches(3.8), Inches(0.9),
                        C_WHITE, "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
        add_text_box(slide, x + Inches(0.2), Inches(6.1), Inches(3.4), Inches(0.3),
                     name, font_size=12, bold=True, color=C_TEXT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), Inches(6.4), Inches(3.4), Inches(0.3),
                     desc, font_size=10, color=C_MUTED, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 18, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 19: Heatmaps
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '谱结构热力图', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Per-Layer Spectral Dynamics', font_size=16, color=C_ACCENT)
    add_figure(slide, 'fig13_heatmaps_v2.png', Inches(0.5), Inches(1.4), width=Inches(12.0), height=Inches(5.5))
    add_page_number(slide, 19, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 20: Section - Part IV
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_DARK)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.4),
                 "PART IV", font_size=14, color=RGBColor(0x6B, 0x7B, 0x8D),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(1.0),
                 "实践价值", font_size=44, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_text_box(slide, Inches(1.5), Inches(4.4), Inches(10.3), Inches(0.5),
                 "Practical Applications for Pretraining Engineers", font_size=18,
                 color=RGBColor(0x8E, 0x9C, 0xAD), alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 20, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 21: Dashboard
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '训练监控仪表盘', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Health Monitoring Dashboard', font_size=15, color=C_ACCENT)
    add_figure(slide, 'fig09_dashboard.png', Inches(0.3), Inches(1.4), height=Inches(5.0))

    # Alert levels
    alerts = [
        ("🟢 HEALTHY", "dα/dt < 0 → 继续训练", C_LIGHT_GREEN, C_EMERALD),
        ("🟡 PLATEAU", "|dα/dt| ≈ 0 → 考虑 decay", C_LIGHT_AMBER, C_AMBER),
        ("🔴 REVERSAL", "dα/dt > 0 ×3 → 立即降 LR", C_LIGHT_RED, C_RED),
        ("⚫ EXHAUSTED", "α > 4 + ↑ → 需要更多容量", RGBColor(0xF2, 0xF4, 0xF7), C_TEXT),
    ]
    for i, (label, desc, bg, tc) in enumerate(alerts):
        y = Inches(1.8 + i * 1.2)
        add_rounded_rect(slide, Inches(7.5), y, Inches(5.0), Inches(0.9), bg, "")
        add_text_box(slide, Inches(7.7), y + Inches(0.1), Inches(4.6), Inches(0.35),
                     label, font_size=13, bold=True, color=tc)
        add_text_box(slide, Inches(7.7), y + Inches(0.45), Inches(4.6), Inches(0.35),
                     desc, font_size=11, color=C_SEC, font_name=FONT_CN)

    add_text_box(slide, Inches(7.5), Inches(6.6), Inches(5.0), Inches(0.4),
                 "成本: < 0.05% 训练 compute", font_size=12, color=C_MUTED, font_name=FONT_CN)
    add_page_number(slide, 21, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 22: Post-hoc Diagnosis
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 'Post-hoc 模型体检', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 '5-Minute Model Audit', font_size=16, color=C_ACCENT)

    # K2 case
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "案例: K2-65B 诊断", font_size=18, bold=True, color=C_TEXT, font_name=FONT_CN)
    k2_data = [["SR/d", "0.036", "偏低"], ["α", "5.09", "不成熟"], ["Reversal", "Δα = +0.65", "已退化"]]
    tbl = slide.shapes.add_table(3, 3, Inches(0.8), Inches(2.0), Inches(5.0), Inches(1.8)).table
    for r in range(3):
        for c in range(3):
            cell = tbl.cell(r, c)
            cell.text = k2_data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.name = FONT_MONO if c == 1 else FONT_MAIN
            if c == 2:
                p.font.color.rgb = C_RED

    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(5.0), Inches(0.8),
                 "诊断: 训练不充分 (only 1.4T tokens, D/N=21)\n外部验证: K2 在 15/21 benchmark 低于 Llama-2-70B",
                 font_size=12, color=C_SEC, font_name=FONT_CN)
    add_rounded_rect(slide, Inches(0.8), Inches(5.3), Inches(5.0), Inches(0.7),
                    C_LIGHT_GREEN, "")
    add_text_box(slide, Inches(1.0), Inches(5.4), Inches(4.6), Inches(0.5),
                 "5 分钟测量 = 跑几天 benchmark 的结论",
                 font_size=14, bold=True, color=C_EMERALD, font_name=FONT_CN)

    # Audit table
    add_text_box(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "完整模型审计表", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    audit = [
        ["Model", "SR/d", "α", "判定"],
        ["Pythia-1B", "0.050", "2.78", "✅ 成熟"],
        ["OLMo-2-1B", "0.064", "2.37", "✅ 成熟"],
        ["OLMo-2-13B", "0.043", "6.95", "❌ 不成熟"],
        ["OLMo-2-32B", "0.043", "5.25", "⚠️ 不成熟"],
        ["K2-65B", "0.036", "5.09", "❌ 严重不足"],
    ]
    tbl = slide.shapes.add_table(6, 4, Inches(6.8), Inches(2.0), Inches(5.5), Inches(3.8)).table
    for r in range(6):
        for c in range(4):
            cell = tbl.cell(r, c)
            cell.text = audit[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.name = FONT_MONO if c in [1, 2] else FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
    add_page_number(slide, 22, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 23: Recommendations
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '对预训练的建议', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Actionable Recommendations', font_size=15, color=C_ACCENT)

    recs = [
        ("建议 1: 不用 Cosine", 'Cosine 过早衰减 LR，浪费"结构形成黄金期"。\n\nWSD 或 α-Guided 保持高 LR 到 74-83%\n→ benchmark +2-2.6%', C_ACCENT),
        ("建议 2: 数据远超 Chinchilla", "结构成熟需要 15-25× Chinchilla 推荐量\n\n1B: ~300B tokens\n7B: ~3.5T tokens\n13B: 远超 5T tokens", C_RED),
        ("建议 3: 关注 MLP 容量", "大模型 Attn 已结构化,\nMLP 是瓶颈 (gap 可达 5.4)\n\n→ 更宽 FFN / MoE 增加 MLP 容量", C_TEAL),
    ]
    for i, (title, body, color) in enumerate(recs):
        x = Inches(0.8 + i * 4.2)
        add_rounded_rect(slide, x, Inches(1.7), Inches(3.8), Inches(4.0),
                        C_WHITE, "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
        # Top color bar
        add_rounded_rect(slide, x, Inches(1.7), Inches(3.8), Inches(0.08), color, "")
        add_text_box(slide, x + Inches(0.2), Inches(1.9), Inches(3.4), Inches(0.5),
                     title, font_size=14, bold=True, color=color, font_name=FONT_CN)
        add_text_box(slide, x + Inches(0.2), Inches(2.5), Inches(3.4), Inches(3.0),
                     body, font_size=12, color=C_SEC, font_name=FONT_CN)

    add_rounded_rect(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(1.0), Inches(6.2), Inches(11.0), Inches(0.6),
                 "总结: 发布前 5 分钟谱体检 — SR/d 接近 0.040 + 0.61/√d 且 α < 4 → 训练充分",
                 font_size=14, bold=True, color=C_ACCENT, font_name=FONT_CN)
    add_page_number(slide, 23, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 24: Comparison Table
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '与现有方法对比', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Comparison with Existing Methods', font_size=14, color=C_ACCENT)

    comp = [
        ["Capability", "Loss", "WeightWatcher", "Scaling Laws", "Ours"],
        ["实时监控", "✓", "✗ (post-hoc)", "✗", "✓"],
        ["跨模型可比", "✗", "✓", "✗", "✓"],
        ["检测结构退化", "✗", "✗", "✗", "✓ (α reversal)"],
        ["预测下游性能", "✗", "部分 (r~0.7)", "仅 loss", "ρ=-0.90, R²=0.84"],
        ["指导训练决策", "✗", "✗", "仅数据量", "✓ (+2.56% @1B)"],
        ["开销", "0", "分钟级", "0", "<0.05% compute"],
    ]
    tbl = slide.shapes.add_table(7, 5, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8)).table
    for r in range(7):
        for c in range(5):
            cell = tbl.cell(r, c)
            cell.text = comp[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = FONT_MAIN
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = C_SEC
            elif c == 4:
                p.font.bold = True
                p.font.color.rgb = C_ACCENT

    add_rounded_rect(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6),
                    C_LIGHT_BLUE, "")
    add_text_box(slide, Inches(1.0), Inches(6.7), Inches(11.0), Inches(0.4),
                 "首次将 post-hoc 谱分析升级为实时训练指导工具，且在预测精度和可操作性上全面超越",
                 font_size=13, bold=True, color=C_ACCENT, font_name=FONT_CN)
    add_page_number(slide, 24, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 25: Theory
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.8),
                 '理论基础', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(6.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Theoretical Foundation', font_size=16, color=C_ACCENT)

    # Left: info theory + SNR
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "信息论恒等式", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_rounded_rect(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.8),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(1.0), Inches(2.1), Inches(5.1), Inches(0.5),
                 "SR(W) = exp(H₂) = 1 / Σ pᵢ²\nwhere pᵢ = σᵢ² / ||W||²F",
                 font_size=13, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(0.4),
                 "α Reversal 的 SNR 解释", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_rounded_rect(slide, Inches(0.8), Inches(3.7), Inches(5.5), Inches(0.8),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(1.0), Inches(3.8), Inches(5.1), Inches(0.5),
                 "SNR(t) = ||E[∇L]||² / Var[∇L]\nα reversal ⟺ SNR(t) < 1",
                 font_size=13, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(0.4),
                 "降 η 使 noise ∝ η² ↓ → 恢复 SNR > 1",
                 font_size=12, color=C_SEC, font_name=FONT_CN)

    # Right: Wedin + Langevin
    add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.4),
                 "Wedin 子空间稳定性", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_rounded_rect(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.7),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.2), Inches(2.1), Inches(5.1), Inches(0.4),
                 "||sinΘ(Uk, U'k)|| ≤ η·||∇L|| / γk(W)",
                 font_size=13, color=C_TEXT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), Inches(2.9), Inches(5.5), Inches(0.4),
                 "α-guided 在 γk 缩小时减小 η → 最小化子空间旋转",
                 font_size=12, color=C_SEC, font_name=FONT_CN)

    add_text_box(slide, Inches(7.0), Inches(3.7), Inches(5.5), Inches(0.4),
                 "Langevin 动力学", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_rounded_rect(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(1.2),
                    RGBColor(0xF2, 0xF4, 0xF7), "")
    add_text_box(slide, Inches(7.2), Inches(4.3), Inches(5.1), Inches(1.0),
                 "SGD = 梯度下降 + 有效温度扩散:\ndσᵢ = -η ∂L/∂σᵢ dt + √(2ηTeff) dBᵢ\nHigher Teff → spectrum blurs → α ↑",
                 font_size=12, color=C_TEXT, font_name=FONT_MONO)
    add_page_number(slide, 25, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 26: Summary
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_BG)
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(8.0), Inches(0.8),
                 '总结与贡献', font_size=32, bold=True, color=C_TEXT, font_name=FONT_CN)
    add_text_box(slide, Inches(8.0), Inches(0.7), Inches(5.0), Inches(0.5),
                 'Summary', font_size=16, color=C_ACCENT)

    # Contributions
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
                 "四个层次的贡献", font_size=16, bold=True, color=C_TEXT, font_name=FONT_CN)
    contribs = [
        "描述性 (Descriptive): SR/d 通用收敛 + α reversal",
        "预测性 (Predictive): 仅从权重预测 84% 下游方差",
        "指导性 (Prescriptive): α-guided +2.56% @ 1B",
        "诊断性 (Diagnostic): 5 min post-hoc 模型审计",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5), contribs, font_size=13)

    # Science significance
    add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.0),
                    RGBColor(0xF2, 0xF4, 0xF7), "", border_color=RGBColor(0xDE, 0xE2, 0xE8))
    add_text_box(slide, Inches(1.0), Inches(4.9), Inches(5.1), Inches(0.3),
                 "科学意义", font_size=11, bold=True, color=C_MUTED, font_name=FONT_CN)
    science = [
        "通用压缩定律: ΔH₂ ≈ -2 nats",
        "结构相变: N ≈ 1.7B 阈值",
        "Loss ≠ Quality 的直接证据",
        "SR = exp(H₂) 严格桥梁",
    ]
    add_bullet_list(slide, Inches(1.0), Inches(5.3), Inches(5.0), Inches(1.4), science, font_size=12)

    # Right: metrics
    metrics_summary = [("14", "Models"), ("4", "Architectures"), ("~340", "Checkpoints"), ("930×", "Scale Range")]
    for i, (val, label) in enumerate(metrics_summary):
        row, col = divmod(i, 2)
        x = Inches(7.0 + col * 2.8)
        y = Inches(1.8 + row * 1.8)
        add_rounded_rect(slide, x, y, Inches(2.5), Inches(1.4),
                        RGBColor(0xF2, 0xF4, 0xF7), "")
        add_text_box(slide, x, y + Inches(0.15), Inches(2.5), Inches(0.7),
                     val, font_size=30, bold=True, color=C_ACCENT, font_name=FONT_MONO, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + Inches(0.8), Inches(2.5), Inches(0.4),
                     label, font_size=12, color=C_MUTED, alignment=PP_ALIGN.CENTER)

    # Target
    add_rounded_rect(slide, Inches(7.0), Inches(5.6), Inches(5.3), Inches(1.2),
                    C_WHITE, "", border_color=C_ACCENT)
    add_text_box(slide, Inches(7.0), Inches(5.75), Inches(5.3), Inches(0.4),
                 "Target: NeurIPS 2026", font_size=16, bold=True, color=C_ACCENT, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.0), Inches(6.2), Inches(5.3), Inches(0.5),
                 "14 publication-quality figures\n38 citations in Related Work",
                 font_size=11, color=C_MUTED, alignment=PP_ALIGN.CENTER)
    add_page_number(slide, 26, TOTAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 27: Thank You
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide, C_DARK)
    add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
                 "Thank You", font_size=48, bold=True, color=C_WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.5),
                 "Beyond Loss Curves: Spectral Monitoring for LLM Pretraining",
                 font_size=16, color=RGBColor(0x8E, 0x9C, 0xAD), alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(1.0),
                 "仅从权重矩阵的谱结构，就能诊断训练健康度、预测下游性能、\n指导 LR 决策——成本 <0.05% 训练 compute",
                 font_size=14, color=RGBColor(0x8E, 0x9C, 0xAD),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_CN)
    add_text_box(slide, Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.4),
                 "Q & A", font_size=14, color=RGBColor(0x6B, 0x7B, 0x8D),
                 alignment=PP_ALIGN.CENTER, font_name=FONT_MONO)
    add_page_number(slide, 27, TOTAL)

    # ─── Save ──────────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'presentation_v3.pptx')
    prs.save(output_path)
    print(f"✓ Saved: {output_path}")
    print(f"  {TOTAL} slides, widescreen 16:9")


if __name__ == '__main__':
    main()
